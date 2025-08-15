import os
from pathlib import Path
import logging
from datetime import datetime

# Document processing libraries
from docx import Document
from docx.shared import Inches
import pypdf
from pypdf import PdfWriter, PdfReader
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches as PptxInches

# Image processing
from PIL import Image, ImageEnhance, ImageFilter
import img2pdf

# Additional utilities
import tempfile
import shutil
import threading

from sklearn.linear_model import enet_path

class FileConverter:
    def __init__(self):
        self.setup_logging()
        self.conversion_callbacks = {}
        
    def setup_logging(self):
        """Setup logging for error tracking"""
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler('converter.log'),
                logging.StreamHandler()
            ]
        )
        self.logger = logging.getLogger(__name__)
        
    def set_progress_callback(self, callback):
        """Set callback for progress updates"""
        self.progress_callback = callback
        
    def update_progress(self, percentage, message=""):
        """Update progress if callback is set"""
        if hasattr(self, 'progress_callback') and self.progress_callback:
            self.progress_callback(percentage, message)
            
    def pdf_to_word(self, input_path, output_dir):
        """Convert PDF to Word document with enhanced error handling"""
        try:
            input_path = Path(input_path)
            output_path = Path(output_dir) / f"{input_path.stem}_converted.docx"
            
            self.logger.info(f"Converting PDF to Word: {input_path}")
            self.update_progress(10, "Initializing PDF conversion...")
            
            # Check if PDF is encrypted
            reader = PdfReader(str(input_path))
            if reader.is_encrypted:
                raise Exception("PDF is password protected. Please unlock it first.")
            
            self.update_progress(30, "Analyzing PDF structure...")
            
            # Use pdf2docx for conversion
            cv = Converter(str(input_path))
            self.update_progress(50, "Converting pages...")
            
            cv.convert(str(output_path), start=0, end=None)
            cv.close()
            
            self.update_progress(90, "Finalizing document...")
            
            # Verify output file exists and has content
            if not output_path.exists() or output_path.stat().st_size < 1000:
                raise Exception("Conversion failed - output file is empty or corrupted")
            
            self.update_progress(100, "PDF to Word conversion completed!")
            self.logger.info(f"Conversion completed: {output_path}")
            return str(output_path)
            
        except Exception as e:
            self.logger.error(f"Error converting PDF to Word: {str(e)}")
            raise Exception(f"PDF to Word conversion failed: {str(e)}")
            
    def word_to_pdf(self, input_path, output_dir):
        """Convert Word document to PDF with multiple fallback methods"""
        try:
            input_path = Path(input_path)
            output_path = Path(output_dir) / f"{input_path.stem}_converted.pdf"
            
            self.logger.info(f"Converting Word to PDF: {input_path}")
            self.update_progress(10, "Loading Word document...")
            
            # Try different methods based on platform and available libraries
            success = False
            
            # Method 1: Windows with Office (best quality)
            if os.name == 'nt':
                success = self._word_to_pdf_win32(input_path, output_path)
                
            # Method 2: LibreOffice (cross-platform)
            if not success:
                success = self._word_to_pdf_libreoffice(input_path, output_path)
                
            # Method 3: python-docx + reportlab (fallback)
            if not success:
                success = self._word_to_pdf_fallback(input_path, output_path)
                
            if not success:
                raise Exception("All conversion methods failed")
                
            self.update_progress(100, "Word to PDF conversion completed!")
            self.logger.info(f"Conversion completed: {output_path}")
            return str(output_path)
            
        except Exception as e:
            self.logger.error(f"Error converting Word to PDF: {str(e)}")
            raise Exception(f"Word to PDF conversion failed: {str(e)}")
            
    def _word_to_pdf_win32(self, input_path, output_path):
        """Convert using Win32 COM (Windows only)"""
        try:
            import win32com.client
            
            self.update_progress(30, "Initializing Word application...")
            
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            
            self.update_progress(50, "Opening document...")
            doc = word.Documents.Open(str(input_path))
            
            self.update_progress(80, "Converting to PDF...")
            doc.SaveAs(str(output_path), FileFormat=17)  # 17 = PDF format
            
            doc.Close()
            word.Quit()
            
            return output_path.exists()
            
        except ImportError:
            self.logger.info("Win32com not available")
            return False
        except Exception as e:
            self.logger.error(f"Win32 conversion failed: {str(e)}")
            return False
            
    def _word_to_pdf_libreoffice(self, input_path, output_path):
        """Convert using LibreOffice (cross-platform)"""
        try:
            import subprocess
            
            self.update_progress(40, "Using LibreOffice for conversion...")
            
            # Try different LibreOffice commands
            libreoffice_commands = [
                'libreoffice',
                'soffice',
                '/Applications/LibreOffice.app/Contents/MacOS/soffice'  # macOS
            ]
            
            for cmd in libreoffice_commands:
                try:
                    result = subprocess.run([
                        cmd, '--headless', '--convert-to', 'pdf',
                        '--outdir', str(output_path.parent), str(input_path)
                    ], capture_output=True, timeout=60)
                    
                    if result.returncode == 0:
                        return output_path.exists()
                        
                except (FileNotFoundError, subprocess.TimeoutExpired):
                    continue
                    
            return False
            
        except Exception as e:
            self.logger.error(f"LibreOffice conversion failed: {str(e)}")
            return False
            
    def _word_to_pdf_fallback(self, input_path, output_path):
        """Fallback conversion using python-docx + reportlab"""
        try:
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            
            self.update_progress(40, "Using fallback conversion method...")
            
            # Read Word document
            doc = Document(str(input_path))
            
            self.update_progress(60, "Processing document content...")
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(str(output_path), pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Add title if available
            if doc.core_properties.title:
                title_style = ParagraphStyle(
                    'CustomTitle',
                    parent=styles['Title'],
                    fontSize=18,
                    spaceAfter=30
                )
                story.append(Paragraph(doc.core_properties.title, title_style))
                story.append(Spacer(1, 12))
            
            # Process paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    # Determine style based on paragraph format
                    if para.style.name.startswith('Heading'):
                        style = styles['Heading1']
                    elif para.style.name == 'Title':
                        style = styles['Title']
                    else:
                        style = styles['Normal']
                    
                    try:
                        p = Paragraph(para.text, style)
                        story.append(p)
                        story.append(Spacer(1, 12))
                    except:
                        # Skip problematic paragraphs
                        continue
            
            self.update_progress(90, "Building PDF...")
            pdf_doc.build(story)
            
            return output_path.exists()
            
        except ImportError:
            self.logger.error("Reportlab not available for fallback conversion")
            return False
        except Exception as e:
            self.logger.error(f"Fallback conversion failed: {str(e)}")
            return False
            
    def word_to_ppt(self, input_path, output_dir):
        """Convert Word document to PowerPoint with smart content detection"""
        try:
            input_path = Path(input_path)
            output_path = Path(output_dir) / f"{input_path.stem}_converted.pptx"
            
            self.logger.info(f"Converting Word to PowerPoint: {input_path}")
            self.update_progress(10, "Loading Word document...")
            
            # Read Word document
            doc = Document(str(input_path))
            
            self.update_progress(30, "Analyzing document structure...")
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Process document with intelligent slide creation
            slides_created = self._create_slides_from_word(doc, prs)
            
            self.update_progress(90, "Saving presentation...")
            prs.save(str(output_path))
            
            self.update_progress(100, f"Created {slides_created} slides successfully!")
            self.logger.info(f"Conversion completed: {output_path}")
            return str(output_path)
            
        except Exception as e:
            self.logger.error(f"Error converting Word to PowerPoint: {str(e)}")
            raise Exception(f"Word to PowerPoint conversion failed: {str(e)}")
            
    def _create_slides_from_word(self, doc, prs):
        """Intelligently create slides from Word document"""
        slide_content = []
        current_title = None
        slides_created = 0
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
                
            # Detect titles/headings for new slides
            is_heading = (
                para.style.name.startswith('Heading') or
                para.style.name == 'Title' or
                (len(text) < 100 and (text.isupper() or text.istitle()))
            )
            
            if is_heading:
                # Create slide with previous content
                if slide_content and current_title:
                    self._create_slide(prs, current_title, slide_content)
                    slides_created += 1
                    slide_content = []
                    
                current_title = text
            else:
                slide_content.append(text)
        
        # Add final slide
        if current_title or slide_content:
            title = current_title or "Content"
            content = slide_content or ["Document content"]
            self._create_slide(prs, title, content)
            slides_created += 1
        
        # If no slides created, create a summary slide
        if slides_created == 0:
            all_content = [p.text for p in doc.paragraphs if p.text.strip()]
            self._create_slide(prs, Path(enet_path).stem, all_content[:10])  # Limit content
            slides_created = 1
            
        return slides_created
        
    def _create_slide(self, presentation, title, content_list):
        """Create a slide with title and content"""
        slide_layout = presentation.slide_layouts[1]  # Title and Content layout
        slide = presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title[:100]  # Limit title length
        
        # Add content
        if slide.placeholders[1]:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()  # Clear existing content
            
            for i, content in enumerate(content_list[:10]):  # Limit to 10 items
                if content.strip():
                    if i == 0:
                        text_frame.text = content[:200]  # Limit content length
                    else:
                        p = text_frame.add_paragraph()
                        p.text = content[:200]
                        
    def ppt_to_word(self, input_path, output_dir):
        """Convert PowerPoint to Word document with enhanced formatting"""
        try:
            input_path = Path(input_path)
            output_path = Path(output_dir) / f"{input_path.stem}_converted.docx"
            
            self.logger.info(f"Converting PowerPoint to Word: {input_path}")
            self.update_progress(10, "Loading PowerPoint presentation...")
            
            # Read PowerPoint presentation
            prs = Presentation(str(input_path))
            
            self.update_progress(30, "Creating Word document...")
            
            # Create Word document
            doc = Document()
            
            # Add title page
            doc.add_heading(f"Converted from {input_path.name}", 0)
            doc.add_paragraph(f"Conversion date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            doc.add_paragraph(f"Total slides: {len(prs.slides)}")
            doc.add_page_break()
            
            self.update_progress(50, "Processing slides...")
            
            for i, slide in enumerate(prs.slides):
                self.update_progress(50 + (30 * i / len(prs.slides)), f"Processing slide {i+1}...")
                
                # Add slide header
                doc.add_heading(f"Slide {i + 1}", level=1)
                
                # Extract text from all shapes
                slide_content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        try:
                            # Clean and sanitize text content
                            text_content = shape.text.strip()
                            # Remove or replace problematic characters
                            text_content = self._sanitize_text_content(text_content)
                            
                            if hasattr(shape, "text_frame") and shape.text_frame:
                                # Handle text frames with multiple paragraphs
                                for paragraph in shape.text_frame.paragraphs:
                                    if paragraph.text.strip():
                                        para_text = self._sanitize_text_content(paragraph.text.strip())
                                        if para_text:
                                            slide_content.append(para_text)
                            else:
                                if text_content:
                                    slide_content.append(text_content)
                        except Exception as e:
                            self.logger.warning(f"Error processing shape text: {str(e)}")
                            continue
                
                # Add content to document
                if slide_content:
                    for content in slide_content:
                        try:
                            # Add paragraph with error handling
                            doc.add_paragraph(content)
                        except Exception as e:
                            self.logger.warning(f"Error adding paragraph: {str(e)}")
                            # Add a safe fallback
                            doc.add_paragraph("[Content could not be processed]")
                else:
                    doc.add_paragraph("[No text content found in this slide]")
                
                # Add separator except for last slide
                if i < len(prs.slides) - 1:
                    doc.add_paragraph("─" * 50)
            
            self.update_progress(90, "Saving Word document...")
            doc.save(str(output_path))
            
            self.update_progress(100, "PowerPoint to Word conversion completed!")
            self.logger.info(f"Conversion completed: {output_path}")
            return str(output_path)
            
        except Exception as e:
            self.logger.error(f"Error converting PowerPoint to Word: {str(e)}")
            raise Exception(f"PowerPoint to Word conversion failed: {str(e)}")
            
    def _sanitize_text_content(self, text):
        """Sanitize text content to remove problematic characters"""
        if not text:
            return ""
            
        try:
            # Remove null bytes and control characters
            text = text.replace('\x00', '')
            
            # Replace problematic Unicode characters
            text = text.encode('utf-8', errors='ignore').decode('utf-8')
            
            # Remove or replace other problematic characters
            problematic_chars = {
                '\x01': '',
                '\x02': '',
                '\x03': '',
                '\x04': '',
                '\x05': '',
                '\x06': '',
                '\x07': '',
                '\x08': '',
                '\x0b': '',
                '\x0c': '',
                '\x0e': '',
                '\x0f': '',
                '\x10': '',
                '\x11': '',
                '\x12': '',
                '\x13': '',
                '\x14': '',
                '\x15': '',
                '\x16': '',
                '\x17': '',
                '\x18': '',
                '\x19': '',
                '\x1a': '',
                '\x1b': '',
                '\x1c': '',
                '\x1d': '',
                '\x1e': '',
                '\x1f': ''
            }
            
            for char, replacement in problematic_chars.items():
                text = text.replace(char, replacement)
            
            # Ensure text is properly formatted
            text = ' '.join(text.split())  # Normalize whitespace
            
            return text
            
        except Exception as e:
            self.logger.warning(f"Error sanitizing text: {str(e)}")
            return "[Text processing error]"
            
    def image_to_pdf(self, input_path, output_dir):
        """Convert image to PDF with quality enhancement"""
        try:
            input_path = Path(input_path)
            output_path = Path(output_dir) / f"{input_path.stem}_converted.pdf"
            
            self.logger.info(f"Converting Image to PDF: {input_path}")
            self.update_progress(10, "Loading image...")
            
            # Open and process image
            with Image.open(str(input_path)) as img:
                self.update_progress(30, "Processing image...")
                
                # Get original format and size
                original_format = img.format
                original_size = img.size
                
                # Convert to RGB if necessary
                if img.mode in ('RGBA', 'LA'):
                    # Handle transparency by adding white background
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'RGBA':
                        background.paste(img, mask=img.split()[-1])
                    else:
                        background.paste(img)
                    img = background
                elif img.mode not in ('RGB', 'L'):
                    img = img.convert('RGB')
                
                self.update_progress(50, "Enhancing image quality...")
                
                # Apply quality enhancements
                img = self._enhance_image_quality(img)
                
                self.update_progress(70, "Creating PDF...")
                
                # Optimize image size for PDF
                max_size = (2000, 2000)  # Maximum dimensions
                if img.size[0] > max_size[0] or img.size[1] > max_size[1]:
                    img.thumbnail(max_size, Image.Resampling.LANCZOS)
                
                # Save as temporary high-quality image
                temp_dir = Path(tempfile.mkdtemp())
                temp_image_path = temp_dir / f"optimized_{input_path.stem}.jpg"
                
                img.save(str(temp_image_path), 'JPEG', quality=95, optimize=True)
                
                self.update_progress(90, "Converting to PDF...")
                
                # Convert to PDF using img2pdf with optimized settings
                with open(str(output_path), "wb") as f:
                    # Create PDF with proper page size
                    pdf_bytes = img2pdf.convert(
                        str(temp_image_path),
                        layout_fun=img2pdf.get_layout_fun((img2pdf.mm_to_pt(210), img2pdf.mm_to_pt(297)))  # A4
                    )
                    f.write(pdf_bytes)
                
                # Clean up temporary file
                temp_image_path.unlink()
                temp_dir.rmdir()
            
            self.update_progress(100, "Image to PDF conversion completed!")
            self.logger.info(f"Conversion completed: {output_path}")
            return str(output_path)
            
        except Exception as e:
            self.logger.error(f"Error converting Image to PDF: {str(e)}")
            raise Exception(f"Image to PDF conversion failed: {str(e)}")
            
    def multi_image_to_pdf(self, image_paths, output_dir):
        """Convert multiple images to a single PDF with enhanced processing"""
        try:
            if not image_paths:
                raise Exception("No images provided for conversion")
                
            # Create output filename based on first image
            first_image = Path(image_paths[0])
            output_path = Path(output_dir) / f"multi_image_converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            
            self.logger.info(f"Converting {len(image_paths)} images to PDF")
            self.update_progress(5, f"Processing {len(image_paths)} images...")
            
            processed_images = []
            temp_dir = Path(tempfile.mkdtemp())
            
            try:
                for i, image_path in enumerate(image_paths):
                    progress = 10 + (70 * i / len(image_paths))
                    self.update_progress(progress, f"Processing image {i+1}/{len(image_paths)}...")
                    
                    # Process each image
                    with Image.open(str(image_path)) as img:
                        # Convert to RGB if necessary
                        if img.mode in ('RGBA', 'LA'):
                            background = Image.new('RGB', img.size, (255, 255, 255))
                            if img.mode == 'RGBA':
                                background.paste(img, mask=img.split()[-1])
                            else:
                                background.paste(img)
                            img = background
                        elif img.mode not in ('RGB', 'L'):
                            img = img.convert('RGB')
                        
                        # Apply quality enhancements
                        img = self._enhance_image_quality(img)
                        
                        # Optimize size
                        max_size = (2000, 2000)
                        if img.size[0] > max_size[0] or img.size[1] > max_size[1]:
                            img.thumbnail(max_size, Image.Resampling.LANCZOS)
                        
                        # Save processed image
                        temp_image_path = temp_dir / f"processed_{i:03d}.jpg"
                        img.save(str(temp_image_path), 'JPEG', quality=95, optimize=True)
                        processed_images.append(str(temp_image_path))
                
                self.update_progress(85, "Creating PDF from processed images...")
                
                # Convert all processed images to PDF
                with open(str(output_path), "wb") as f:
                    pdf_bytes = img2pdf.convert(
                        processed_images,
                        layout_fun=img2pdf.get_layout_fun((img2pdf.mm_to_pt(210), img2pdf.mm_to_pt(297)))
                    )
                    f.write(pdf_bytes)
                
                self.update_progress(100, f"Multi-image PDF created with {len(image_paths)} images!")
                
            finally:
                # Clean up temporary files
                for temp_file in processed_images:
                    try:
                        Path(temp_file).unlink()
                    except:
                        pass
                try:
                    temp_dir.rmdir()
                except:
                    pass
            
            self.logger.info(f"Multi-image conversion completed: {output_path}")
            return str(output_path)
            
        except Exception as e:
            self.logger.error(f"Error converting multiple images to PDF: {str(e)}")
            raise Exception(f"Multi-image to PDF conversion failed: {str(e)}")
            
    def _enhance_image_quality(self, img):
        """Enhance image quality for better PDF output"""
        try:
            # Sharpen the image slightly
            enhancer = ImageEnhance.Sharpness(img)
            img = enhancer.enhance(1.1)
            
            # Enhance contrast slightly
            enhancer = ImageEnhance.Contrast(img)
            img = enhancer.enhance(1.05)
            
            # Enhance brightness if too dark
            enhancer = ImageEnhance.Brightness(img)
            img = enhancer.enhance(1.02)
            
            return img
            
        except Exception:
            # Return original if enhancement fails
            return img
            
    def get_file_info(self, file_path):
        """Get comprehensive file information"""
        file_path = Path(file_path)
        
        info = {
            'name': file_path.name,
            'size': file_path.stat().st_size,
            'extension': file_path.suffix,
            'modified': datetime.fromtimestamp(file_path.stat().st_mtime),
            'type': 'Unknown'
        }
        
        # Determine file type and get specific information
        try:
            if file_path.suffix.lower() == '.pdf':
                from pypdf import PdfReader
                reader = PdfReader(str(file_path))
                info.update({
                    'type': 'PDF Document',
                    'pages': len(reader.pages),
                    'encrypted': reader.is_encrypted,
                    'metadata': reader.metadata if reader.metadata else {}
                })
                
            elif file_path.suffix.lower() in ['.docx', '.doc']:
                doc = Document(str(file_path))
                info.update({
                    'type': 'Word Document',
                    'paragraphs': len([p for p in doc.paragraphs if p.text.strip()]),
                    'core_properties': {
                        'title': doc.core_properties.title,
                        'author': doc.core_properties.author,
                        'subject': doc.core_properties.subject
                    }
                })
                
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                prs = Presentation(str(file_path))
                info.update({
                    'type': 'PowerPoint Presentation',
                    'slides': len(prs.slides)
                })
                
            elif file_path.suffix.lower() in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']:
                with Image.open(str(file_path)) as img:
                    info.update({
                        'type': 'Image File',
                        'dimensions': f"{img.width} x {img.height}",
                        'mode': img.mode,
                        'format': img.format
                    })
                    
        except Exception as e:
            self.logger.warning(f"Could not get detailed info for {file_path}: {str(e)}")
        
        return info
        
    def batch_convert(self, file_list, conversion_type, output_dir, progress_callback=None):
        """Convert multiple files with progress tracking"""
        results = []
        failed_files = []
        total_files = len(file_list)
        
        for i, file_path in enumerate(file_list):
            try:
                if progress_callback:
                    progress_callback(
                        (i / total_files) * 100, 
                        f"Converting {Path(file_path).name}... ({i+1}/{total_files})"
                    )
                
                if conversion_type == "pdf_to_word":
                    result = self.pdf_to_word(file_path, output_dir)
                elif conversion_type == "word_to_pdf":
                    result = self.word_to_pdf(file_path, output_dir)
                elif conversion_type == "word_to_ppt":
                    result = self.word_to_ppt(file_path, output_dir)
                elif conversion_type == "ppt_to_word":
                    result = self.ppt_to_word(file_path, output_dir)
                elif conversion_type == "image_to_pdf":
                    result = self.image_to_pdf(file_path, output_dir)
                else:
                    raise Exception(f"Unsupported conversion type: {conversion_type}")
                
                results.append(result)
                
            except Exception as e:
                failed_files.append((file_path, str(e)))
                self.logger.error(f"Failed to convert {file_path}: {str(e)}")
        
        if progress_callback:
            progress_callback(100, f"Batch conversion completed! {len(results)} successful, {len(failed_files)} failed")
            
        return results, failed_files